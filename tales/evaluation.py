"""
Evaluation module for RAG and PowerPoint generation agents.
This module contains functions and classes for evaluating the performance
of the RAG agent and PowerPoint generation agent.
"""
import json
import time
import asyncio
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass, asdict

from langchain_core.messages import AnyMessage, HumanMessage, AIMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.runnables import RunnableConfig
from pptx import Presentation
import numpy as np

from tales.agent import agent, GraphState
from tales.ppt_agent import ppt_agent
from tales.db_handler import ChromaDBHandler
from tales.config import DB_PATH, USE_REASONING, LOCAL_LLMS, llm

# Define evaluation metrics for RAG agent
@dataclass
class RAGMetrics:
    """Metrics for evaluating RAG agent performance."""
    query: str
    response_time: float  # In seconds
    context_relevance_score: float  # 0-10 score
    answer_correctness_score: float  # 0-10 score
    answer_completeness_score: float  # 0-10 score
    hallucination_score: float  # 0-10 score (lower is better)
    num_documents_retrieved: int
    num_research_iterations: int
    
    def to_dict(self):
        """Convert metrics to dictionary."""
        return asdict(self)

# Define evaluation metrics for PowerPoint agent
@dataclass
class PPTMetrics:
    """Metrics for evaluating PowerPoint generation."""
    generation_time: float  # In seconds
    slides_count: int
    avg_content_per_slide: float
    content_coverage_score: float  # 0-10 score
    design_quality_score: float  # 0-10 score
    organization_score: float  # 0-10 score
    
    def to_dict(self):
        """Convert metrics to dictionary."""
        return asdict(self)

class RAGEvaluator:
    """Evaluator for RAG agent performance."""
    
    def __init__(self, db_handler: ChromaDBHandler = None):
        """Initialize RAG evaluator.
        
        Args:
            db_handler: ChromaDB handler for accessing conversations
        """
        self.db_handler = db_handler or ChromaDBHandler(persist_directory=DB_PATH)
        
        # Initialize evaluation LLM (using Gemini for evaluation)
        self.eval_llm = llm

    def evaluate_response(self, query: str, messages: List[AnyMessage], 
                          context_docs: List[Any], iterations: int) -> RAGMetrics:
        """Evaluate the RAG agent's response.
        
        Args:
            query: Original user query
            messages: List of messages including the response
            context_docs: Documents retrieved as context
            iterations: Number of research iterations
            
        Returns:
            RAGMetrics object with evaluation scores
        """
        # Extract the response (last AI message)
        response = next((msg.content for msg in reversed(messages) 
                         if isinstance(msg, AIMessage)), "")
        
        # Get context text from documents
        context_texts = [doc.page_content for doc in context_docs] if context_docs else []
        
        # Evaluate context relevance
        context_relevance = self._evaluate_context_relevance(query, context_texts)
        
        # Evaluate answer quality
        answer_correctness = self._evaluate_answer_correctness(query, response, context_texts)
        
        # Evaluate answer completeness
        answer_completeness = self._evaluate_answer_completeness(query, response)
        
        # Evaluate for hallucinations
        hallucination_score = self._evaluate_hallucinations(response, context_texts)
        
        return RAGMetrics(
            query=query,
            response_time=0.0,  # Will be set by the evaluate_rag_query method
            context_relevance_score=context_relevance,
            answer_correctness_score=answer_correctness,
            answer_completeness_score=answer_completeness,
            hallucination_score=hallucination_score,
            num_documents_retrieved=len(context_docs),
            num_research_iterations=iterations
        )
    
    def evaluate_rag_query(self, query: str, thread_id: int = None) -> Tuple[RAGMetrics, List[AnyMessage]]:
        """Run and evaluate a query through the RAG agent.
        
        Args:
            query: User query to evaluate
            thread_id: Optional thread ID for conversation context
            
        Returns:
            Tuple of (RAGMetrics, messages)
        """
        # Create a new thread ID if not provided
        thread_id = thread_id or int(time.time())
        
        # Get existing messages or start fresh
        messages = self.db_handler.get_conversation(thread_id) if thread_id else []
        
        # Add the query
        messages.append(HumanMessage(content=query))
        
        # Track iterations
        iterations_counter = [0]
        context_docs = []
        
        # Create a modified state graph to track iterations and documents
        def _intercept_retrieve(state: GraphState):
            iterations_counter[0] += 1
            retrieved_docs = state.get("context", [])
            if retrieved_docs:
                context_docs.extend(retrieved_docs)
            return state
        
        # Start timing
        start_time = time.time()
        
        # Initialize state with the query
        initial_state = {
            "context": [],
            "messages": messages,
            "summary": "",
            "query": None,
            "more_research": False
        }
        
        # Run the graph with tracking
        config = {"configurable": {"thread_id": thread_id}}
        result = agent.invoke(initial_state, config)
        
        # Calculate time
        end_time = time.time()
        response_time = end_time - start_time
        
        # Evaluate the results
        metrics = self.evaluate_response(
            query=query,
            messages=result["messages"],
            context_docs=context_docs,
            iterations=iterations_counter[0]
        )
        
        # Update the response time
        metrics.response_time = response_time
        
        return metrics, result["messages"]
    
    def _evaluate_context_relevance(self, query: str, context_texts: List[str]) -> float:
        """Evaluate the relevance of retrieved context to the query.
        
        Args:
            query: The user query
            context_texts: List of context document texts
            
        Returns:
            Score from 0-10 on context relevance
        """
        if not context_texts:
            return 0.0
            
        # Limit context size for evaluation
        combined_context = "\n\n".join(context_texts[:5])  # Limit to first 5 docs
        if len(combined_context) > 8000:
            combined_context = combined_context[:8000] + "..."
            
        eval_prompt = [
            HumanMessage(content=f"""You are an expert evaluator for RAG (Retrieval-Augmented Generation) systems.
            
For the following user query, evaluate how relevant the retrieved context is on a scale from 0 to 10, where:
- 0: Completely irrelevant
- 5: Somewhat relevant but missing key information
- 10: Highly relevant and contains all needed information

Query: "{query}"

Retrieved Context:
{combined_context}

Provide your rating as a single number between 0-10 without explanation or other text.
""")
        ]
        
        try:
            response = self.eval_llm.invoke(eval_prompt).content
            # Extract the numeric score
            score = float(response.strip())
            return min(max(score, 0.0), 10.0)  # Ensure it's between 0-10
        except:
            # Default score if evaluation fails
            return 5.0
    
    def _evaluate_answer_correctness(self, query: str, response: str, context_texts: List[str]) -> float:
        """Evaluate the correctness of the answer based on the context.
        
        Args:
            query: The user query
            response: The agent's response
            context_texts: List of context document texts
            
        Returns:
            Score from 0-10 on answer correctness
        """
        if not context_texts:
            return 5.0  # Middle score if no context
            
        # Limit context size for evaluation
        combined_context = "\n\n".join(context_texts[:3])  # First 3 docs
        if len(combined_context) > 6000:
            combined_context = combined_context[:6000] + "..."
            
        eval_prompt = [
            HumanMessage(content=f"""You are an expert evaluator for RAG (Retrieval-Augmented Generation) systems.
            
For the following user query and context, evaluate how factually correct the response is on a scale from 0 to 10, where:
- 0: Completely incorrect, contains false information
- 5: Mixture of correct and incorrect information
- 10: Completely factually correct based on the context

Query: "{query}"

Context (ground truth):
{combined_context}

Response to evaluate:
{response}

Provide your rating as a single number between 0-10 without explanation or other text.
""")
        ]
        
        try:
            response = self.eval_llm.invoke(eval_prompt).content
            # Extract the numeric score
            score = float(response.strip())
            return min(max(score, 0.0), 10.0)  # Ensure it's between 0-10
        except:
            # Default score if evaluation fails
            return 5.0
    
    def _evaluate_answer_completeness(self, query: str, response: str) -> float:
        """Evaluate the completeness of the answer.
        
        Args:
            query: The user query
            response: The agent's response
            
        Returns:
            Score from 0-10 on answer completeness
        """
        eval_prompt = [
            HumanMessage(content=f"""You are an expert evaluator for question answering systems.
            
For the following user query, evaluate how completely the response answers the question on a scale from 0 to 10, where:
- 0: Does not answer the question at all
- 5: Partially answers the question but misses important aspects
- 10: Fully and comprehensively answers the question

Query: "{query}"

Response:
{response}

Provide your rating as a single number between 0-10 without explanation or other text.
""")
        ]
        
        try:
            response = self.eval_llm.invoke(eval_prompt).content
            # Extract the numeric score
            score = float(response.strip())
            return min(max(score, 0.0), 10.0)  # Ensure it's between 0-10
        except:
            # Default score if evaluation fails
            return 5.0
    
    def _evaluate_hallucinations(self, response: str, context_texts: List[str]) -> float:
        """Evaluate the response for hallucinations (information not in context).
        
        Args:
            response: The agent's response
            context_texts: List of context document texts
            
        Returns:
            Score from 0-10 on hallucination (lower is better)
        """
        if not context_texts:
            return 5.0  # Middle score if no context
            
        # Limit context size for evaluation
        combined_context = "\n\n".join(context_texts[:3])  # First 3 docs
        if len(combined_context) > 6000:
            combined_context = combined_context[:6000] + "..."
            
        eval_prompt = [
            HumanMessage(content=f"""You are an expert evaluator for RAG (Retrieval-Augmented Generation) systems.
            
For the following response and context, evaluate the degree of hallucination on a scale from 0 to 10, where:
- 0: No hallucination, response only contains information from the context
- 5: Some hallucination, response includes some information not in the context
- 10: Complete hallucination, response is unrelated to or contradicts the context

Context (ground truth):
{combined_context}

Response to evaluate:
{response}

Provide your rating as a single number between 0-10 without explanation or other text.
""")
        ]
        
        try:
            response = self.eval_llm.invoke(eval_prompt).content
            # Extract the numeric score
            score = float(response.strip())
            return min(max(score, 0.0), 10.0)  # Ensure it's between 0-10
        except:
            # Default score if evaluation fails
            return 5.0


class PowerPointEvaluator:
    """Evaluator for PowerPoint generation."""
    
    def __init__(self):
        """Initialize PowerPoint evaluator."""
        # Initialize evaluation LLM (using Gemini for evaluation)
        self.eval_llm = llm
    
    async def evaluate_ppt_generation(self, messages: List[AnyMessage], 
                                ppt_path: str = "presentation.pptx") -> PPTMetrics:
        """Generate and evaluate a PowerPoint presentation.
        
        Args:
            messages: The conversation messages
            ppt_path: Path where the PowerPoint file will be saved
            
        Returns:
            PPTMetrics object with evaluation scores
        """
        # Start timing
        start_time = time.time()
        
        # Generate the presentation
        await ppt_agent(messages)
        
        # End timing
        end_time = time.time()
        generation_time = end_time - start_time
        
        # Analyze the generated presentation
        ppt_metrics = self._analyze_presentation(ppt_path, messages)
        
        # Update generation time
        ppt_metrics.generation_time = generation_time
        
        return ppt_metrics
    
    def _analyze_presentation(self, ppt_path: str, messages: List[AnyMessage]) -> PPTMetrics:
        """Analyze the generated PowerPoint presentation.
        
        Args:
            ppt_path: Path to the generated PowerPoint file
            messages: The conversation messages that were used to generate the presentation
            
        Returns:
            PPTMetrics object with analysis results
        """
        try:
            # Load the presentation
            prs = Presentation(ppt_path)
            
            # Count slides
            slides_count = len(prs.slides)
            
            # Calculate average content per slide
            total_content = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        total_content += len(shape.text)
            
            avg_content_per_slide = total_content / slides_count if slides_count > 0 else 0
            
            # Extract conversation content for comparison
            conversation_text = "\n".join([msg.content for msg in messages])
            
            # Evaluate content coverage
            content_coverage = self._evaluate_content_coverage(conversation_text, prs)
            
            # Evaluate design quality
            design_quality = self._evaluate_design_quality(prs)
            
            # Evaluate organization
            organization_score = self._evaluate_organization(prs)
            
            return PPTMetrics(
                generation_time=0.0,  # Will be set by the calling method
                slides_count=slides_count,
                avg_content_per_slide=avg_content_per_slide,
                content_coverage_score=content_coverage,
                design_quality_score=design_quality,
                organization_score=organization_score
            )
        except Exception as e:
            print(f"Error analyzing presentation: {e}")
            # Return default metrics in case of error
            return PPTMetrics(
                generation_time=0.0,
                slides_count=0,
                avg_content_per_slide=0.0,
                content_coverage_score=0.0,
                design_quality_score=0.0,
                organization_score=0.0
            )
    
    def _evaluate_content_coverage(self, conversation_text: str, prs: Presentation) -> float:
        """Evaluate how well the presentation covers the content from the conversation.
        
        Args:
            conversation_text: The text from the conversation
            prs: The Presentation object
            
        Returns:
            Score from 0-10 on content coverage
        """
        # Extract all text from the presentation
        presentation_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    presentation_text += shape.text + "\n"
        
        # Limit text sizes for evaluation
        if len(conversation_text) > 6000:
            conversation_text = conversation_text[:6000] + "..."
        if len(presentation_text) > 6000:
            presentation_text = presentation_text[:6000] + "..."
            
        eval_prompt = [
            HumanMessage(content=f"""You are an expert evaluator for PowerPoint presentations.
            
Evaluate how well the presentation covers the main topics and information from the conversation on a scale from 0 to 10, where:
- 0: None of the important information is covered
- 5: Some important information is covered, but key points are missing
- 10: All important information is comprehensively covered

Conversation:
{conversation_text}

Presentation Text:
{presentation_text}

Provide your rating as a single number between 0-10 without explanation or other text.
""")
        ]
        
        try:
            response = self.eval_llm.invoke(eval_prompt).content
            # Extract the numeric score
            score = float(response.strip())
            return min(max(score, 0.0), 10.0)  # Ensure it's between 0-10
        except:
            # Default score if evaluation fails
            return 5.0
    
    def _evaluate_design_quality(self, prs: Presentation) -> float:
        """Evaluate the design quality of the presentation.
        
        Args:
            prs: The Presentation object
            
        Returns:
            Score from 0-10 on design quality
        """
        # Count various design elements
        total_shapes = 0
        total_charts = 0
        total_images = 0
        total_tables = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                # TODO: Add more specific type checking when python-pptx is available
        
        # Simple heuristic: more varied elements = better design
        # This is a very basic metric and should be enhanced with actual design principles
        design_score = min(10, (total_shapes + total_charts + total_images + total_tables) / len(prs.slides))
        
        return design_score
    
    def _evaluate_organization(self, prs: Presentation) -> float:
        """Evaluate the organization and structure of the presentation.
        
        Args:
            prs: The Presentation object
            
        Returns:
            Score from 0-10 on organization
        """
        # Check if the presentation has a title slide
        has_title_slide = False
        if prs.slides:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if hasattr(shape, "text") and len(shape.text) > 0:
                    has_title_slide = True
                    break
        
        # Check if the presentation has a consistent structure
        slide_texts = []
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            slide_texts.append(slide_text)
        
        # Basic organization score (can be improved with more sophisticated analysis)
        organization_score = 5.0  # Default middle score
        
        # Bonus for having a title slide
        if has_title_slide:
            organization_score += 2.0
            
        # Bonus for having enough slides (more than 3)
        if len(prs.slides) > 3:
            organization_score += 1.0
            
        # Bonus for consistent content length across slides
        if slide_texts:
            text_lengths = [len(text) for text in slide_texts]
            variation = np.std(text_lengths) / max(np.mean(text_lengths), 1)
            if variation < 0.5:  # Low variation means consistent content
                organization_score += 2.0
                
        return min(max(organization_score, 0.0), 10.0)  # Ensure it's between 0-10


def run_batch_evaluation(queries: List[str], save_path: str = "evaluation_results.json") -> Dict:
    """Run batch evaluation on a list of queries.
    
    Args:
        queries: List of queries to evaluate
        save_path: Path to save the evaluation results
        
    Returns:
        Dictionary with evaluation results
    """
    evaluator = RAGEvaluator()
    ppt_evaluator = PowerPointEvaluator()
    results = []
    
    for i, query in enumerate(queries):
        print(f"\nEvaluating query {i+1}/{len(queries)}: {query}")
        
        # Evaluate RAG
        rag_metrics, messages = evaluator.evaluate_rag_query(query)
        
        # Evaluate PowerPoint generation
        try:
            ppt_metrics = asyncio.run(ppt_evaluator.evaluate_ppt_generation(messages))
        except Exception as e:
            print(f"Error evaluating PowerPoint: {e}")
            ppt_metrics = PPTMetrics(
                generation_time=0.0,
                slides_count=0,
                avg_content_per_slide=0.0,
                content_coverage_score=0.0,
                design_quality_score=0.0,
                organization_score=0.0
            )
        
        # Store results
        result = {
            "query": query,
            "rag_metrics": rag_metrics.to_dict(),
            "ppt_metrics": ppt_metrics.to_dict()
        }
        results.append(result)
        
        # Save intermediate results
        with open(save_path, "w") as f:
            json.dump(results, f, indent=2)
        
    return results
